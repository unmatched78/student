import streamlit as st  
from PyPDF2 import PdfReader  
from docx import Document  
from mistralai import Mistral  
from langchain.text_splitter import RecursiveCharacterTextSplitter  
import re  
from pptx import Presentation  
from dotenv import load_dotenv  
import os  
import json  

load_dotenv()  

api_keys = [os.getenv("mistral_api_key1"), os.getenv("mistral_api_key2"), os.getenv("mistral_api_key7")]    
current_key_index = 0  

def get_next_api_key():  
    global current_key_index  
    api_key = api_keys[current_key_index]  
    current_key_index = (current_key_index + 1) % len(api_keys)  
    return api_key  

def create_client():  
    api_key = get_next_api_key()  
    return Mistral(api_key=api_key)  

client = create_client()  
#functions for file extraction, question generation, and evaluation...  
def extract_text_from_file(uploaded_file):  
    file_type = uploaded_file.name.split('.')[-1].lower()  
    if file_type == "pdf":  
        reader = PdfReader(uploaded_file)  
        return "\n".join(page.extract_text() for page in reader.pages if page.extract_text())  
    elif file_type == "docx":  
        doc = Document(uploaded_file)  
        return "\n".join(paragraph.text for paragraph in doc.paragraphs)  
    elif file_type == "txt":  
        return uploaded_file.read().decode("utf-8")  
    elif file_type == "pptx":  
        prs = Presentation(uploaded_file)  
        text = []  
        for slide in prs.slides:  
            for shape in slide.shapes:  
                if hasattr(shape, "text"):  
                    text.append(shape.text)  
        return "\n".join(text)  
    else:  
        return None  

def split_text_into_chunks(text, chunk_size=1000, chunk_overlap=100):  
    splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)  
    return splitter.split_text(text)  

def generate_questions_from_text(text_chunks, prompt, num_questions):  
    combined_text = "\n".join(text_chunks)  
    combined_prompt = (  
        f"{prompt}\nContent:\n{combined_text}\n\n"  
        f"Generate exactly {num_questions} multiple-choice questions with correct answers in JSON format:\n"  
        '[{"question": "string", "options": ["A", "B", "C", "D"], "correct_choice": "string"}, please keep in mind that the correct_choice, will not be a letter only , it will be a complete set; i.e :D. It forms the inner layer of the bone matrix., etc ]'  
    )  

    response = client.chat.complete(  
        model="mistral-large-latest",  
        messages=[  
            {"role": "system", "content": "You are a student question generator that creates multiple-choice questions."},  
            {"role": "user", "content": combined_prompt},  
        ],  
        response_format={"type": "json_object"},  
    )  

    try:  
        return json.loads(response.choices[0].message.content)  
    except json.JSONDecodeError:  
        return []  

def evaluate_answers(questions, student_answers):  
    feedback = []  
    total_correct = 0  

    for i, question in enumerate(questions):  
        options = question.get("options", [])  
        correct_answer_text = question.get("correct_choice", "").strip()  
        correct_letter = ""  
   
        for idx, option in enumerate(options):  
            if option.strip().lower() == correct_answer_text.lower():  
                correct_letter = chr(65 + idx)  
                break  

        correct_answer_display = f"{correct_letter}. {correct_answer_text}"  
   
        student_answer = student_answers.get(f"Q{i+1}", "").strip()  
        student_choice_match = re.match(r"^[A-D]", student_answer)  
        student_choice = student_choice_match.group(0) if student_choice_match else ""  
  
        is_correct = correct_letter == student_choice  
        if is_correct:  
            total_correct += 1  

        feedback.append({  
            "question": question["question"],  
            "options": options,  
            "correct_answer": correct_answer_display,  
            "student_answer": student_choice,  
            "is_correct": is_correct  
        })  

    return feedback, total_correct 
def ask_chatbot(question):  
    api_key = api_keys[2]  
    client = Mistral(api_key=api_key)  
    response = client.chat.complete(  
        model="mistral-large-latest",  
        messages=[  
            {"role": "user", "content": question},  
        ],  
        # response_format={"type": "json_object"},  
    )  
    return response.choices[0].message.content  

st.title("MCQs Quiz Generator and Chatbot")  
uploaded_file = st.file_uploader("Upload a document (PDF, DOCX, TXT, PPTX)", type=["pdf", "docx", "txt", "pptx"])  

if "questions" not in st.session_state:  
    st.session_state.questions = []  
if "student_answers" not in st.session_state:  
    st.session_state.student_answers = {}  

# Main logic for quiz generation...  

if uploaded_file:  
    # Continue with existing file processing and question generation...  
    file_content = extract_text_from_file(uploaded_file)  
    if file_content:  
        text_chunks = split_text_into_chunks(file_content)  
        prompt = st.text_input("Enter your prompt(Optional):")  
        num_questions = st.number_input("Number of questions", min_value=1, max_value=40, step=1)  

        if st.button("Generate Questions"):  
            st.session_state.questions = generate_questions_from_text(text_chunks, prompt, num_questions)  
            st.session_state.student_answers = {}  

        if st.session_state.questions:  
            with st.form("quiz_form"):  
                for i, q in enumerate(st.session_state.questions, start=1):  
                    st.subheader(f"Q{i}: {q['question']}")  
                    st.session_state.student_answers[f"Q{i}"] = st.radio(  
                        f"Select an answer for Q{i}",  
                        q['options'],  
                        key=f"q{i}",  
                        index=q['options'].index(st.session_state.student_answers.get(f"Q{i}", q['options'][0]))  
                    )  

                submit_button = st.form_submit_button("Submit Answers")  

            if submit_button:  
                feedback, total_marks = evaluate_answers(st.session_state.questions, st.session_state.student_answers)  
                st.write(f"Your total marks: {total_marks}/{len(st.session_state.questions)}")  
                for f in feedback:
                                        # Access options directly from feedback item  
                    options = f['options']  
                    student_option_index = ord(f['student_answer']) - 65 if f['student_answer'] else None  
                    student_option_text = options[student_option_index] if student_option_index is not None and 0 <= student_option_index < len(options) else 'None'  

                    st.write(f"**Question:** {f['question']}")  
                    st.write(f"**Your Answer:** {f['student_answer']}. {student_option_text}")  
                    st.write(f"**Correct Answer:** {f['correct_answer']}")  
                    st.write("✅ Correct" if f['is_correct'] else "❌ Incorrect")  
        else:  
            st.info("Click 'Generate Questions' to get started.")  
    else:  
        st.error("Unsupported file format or empty document.")

# Chatbot functionality  
st.sidebar.title("Chatbot")  
if st.sidebar.button("Open Chatbot"):  
    st.session_state.chat_open = not st.session_state.get("chat_open", False)  

if st.session_state.get("chat_open", False):  
    st.sidebar.subheader("Ask your question:")  
    user_question = st.sidebar.text_area("Enter your question here:", height=150)  

    if st.sidebar.button("Send Question"):  
        if user_question:  
            chatbot_response = ask_chatbot(user_question)  
            st.sidebar.write("**Chatbot Response:**")  
            st.sidebar.write(chatbot_response)  
        else:  
            st.sidebar.warning("Please enter a question before sending.")